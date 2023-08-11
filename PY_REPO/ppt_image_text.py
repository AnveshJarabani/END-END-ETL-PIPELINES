from pptx import Presentation
import pytesseract
from PIL import Image
import io

# Load the PowerPoint presentation
pptx_path = '../../OMS-CY-205248_4.pptx'
presentation = Presentation(pptx_path)

# Initialize Tesseract
pytesseract.pytesseract.tesseract_cmd = r'C:\Users\ajarabani\AppData\Local\Programs\Python\Python311\Scripts\pytesseract.exe'  # Update this path

# Function to extract text from an image using pytesseract
def extract_text_from_image(image_bytes,lang='eng'):
    image = Image.open(io.BytesIO(image_bytes))
    text = pytesseract.image_to_string(image)
    return text

# Iterate through slides
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Shape type for images
            image_part = shape.image
            if image_part is not None:
                image_bytes = image_part.blob
                # with open('temp_image.png', 'wb') as temp_image_file:
                #     temp_image_file.write(image_bytes)

                extracted_text = extract_text_from_image(image_bytes)
                print("Extracted Text:", extracted_text)